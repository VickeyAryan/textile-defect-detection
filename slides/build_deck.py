"""Build the 'Computer Vision: A Crash Course' slide deck.

Generates a .pptx that teaches CV fundamentals to postgraduate students with
no prior background, using the Fabric Defect Detector project (this repo) as
the running case study. Visible slide text stays concise; the depth lives in
speaker notes, written as a lecturer's narration for each slide.
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

ASSETS = os.path.join(os.path.dirname(__file__), "assets")
MODELS = os.path.join(os.path.dirname(__file__), "..", "models")
OUT_PATH = os.path.join(os.path.dirname(__file__), "Computer_Vision_Crash_Course.pptx")

FONT = "Segoe UI"

def C(hex_str):
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

BLUE, AQUA, YELLOW, GREEN, VIOLET, RED, MAGENTA, ORANGE = (
    C("#2a78d6"), C("#1baf7a"), C("#eda100"), C("#008300"),
    C("#4a3aa7"), C("#e34948"), C("#e87ba4"), C("#eb6834"),
)
SURFACE = C("#fcfcfb")
PAGE = C("#f9f9f7")
TEXT_PRIMARY = C("#0b0b0b")
TEXT_SECONDARY = C("#52514e")
MUTED = C("#898781")
GRIDLINE = C("#e1e0d9")
WHITE = C("#ffffff")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

section_counter = {"n": 0}
current_section = {"label": ""}


# ---------------------------------------------------------------------------
# low-level helpers
# ---------------------------------------------------------------------------
def new_slide():
    return prs.slides.add_slide(BLANK)


def fill_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = MUTED
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, text, left, top, width, height, size=18, color=TEXT_PRIMARY,
             bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT, line_spacing=1.0, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_bullets(slide, items, left, top, width, height, size=17,
                 color=TEXT_PRIMARY, space_after=10, font=FONT, marker=True):
    """items: list of (text, level) tuples, level 0 or 1."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.08
        indent = "      " if level == 1 else ""
        bullet_char = ("–  " if level == 1 else "•  ") if marker else ""
        run = p.add_run()
        run.text = f"{indent}{bullet_char}{text}"
        run.font.size = Pt(size - 2 if level == 1 else size)
        run.font.color.rgb = color if level == 0 else TEXT_SECONDARY
        run.font.name = font
    return box


def add_picture_contain(slide, image_path, left, top, max_w, max_h, align="center"):
    with Image.open(image_path) as im:
        iw, ih = im.size
    ratio = min(max_w / iw, max_h / ih)
    w, h = int(iw * ratio), int(ih * ratio)
    if align == "center":
        l = left + (max_w - w) // 2
    else:
        l = left
    t = top + (max_h - h) // 2
    slide.shapes.add_picture(image_path, l, t, width=w, height=h)


def add_footer(slide, label):
    add_text(slide, label, Inches(0.5), Inches(7.13), Inches(6), Inches(0.3),
              size=10, color=MUTED, font=FONT)
    section_counter["n"] += 1
    add_text(slide, f"{section_counter['n']}", Inches(12.6), Inches(7.13), Inches(0.5), Inches(0.3),
              size=10, color=MUTED, align=PP_ALIGN.RIGHT, font=FONT)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_title_bar(slide, title, accent=BLUE):
    add_rect(slide, Inches(0.5), Inches(0.45), Inches(0.09), Inches(0.62), accent)
    add_text(slide, title, Inches(0.72), Inches(0.4), Inches(12), Inches(0.75),
              size=28, bold=True, color=TEXT_PRIMARY)


# ---------------------------------------------------------------------------
# slide builders
# ---------------------------------------------------------------------------
def title_slide(title, subtitle, tag):
    s = new_slide()
    fill_bg(s, BLUE)
    add_rect(s, Inches(0), Inches(6.55), SLIDE_W, Inches(0.95), C("#184f95"))
    add_text(s, title, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.6),
              size=44, bold=True, color=WHITE)
    add_text(s, subtitle, Inches(0.9), Inches(3.75), Inches(11.5), Inches(0.9),
              size=20, color=C("#cde2fb"))
    add_text(s, tag, Inches(0.9), Inches(6.75), Inches(11), Inches(0.5),
              size=13, color=C("#cde2fb"))
    return s


def section_slide(number, title, description):
    s = new_slide()
    fill_bg(s, TEXT_PRIMARY)
    add_text(s, f"PART {number}", Inches(0.9), Inches(2.7), Inches(4), Inches(0.6),
              size=16, bold=True, color=C("#86b6ef"))
    add_text(s, title, Inches(0.9), Inches(3.15), Inches(11.3), Inches(1.6),
              size=38, bold=True, color=WHITE)
    add_text(s, description, Inches(0.9), Inches(4.35), Inches(10.8), Inches(1.0),
              size=16, color=C("#c3c2b7"))
    current_section["label"] = title
    return s


def content_slide(title, footer_note=None):
    s = new_slide()
    fill_bg(s, SURFACE)
    add_title_bar(s, title)
    add_footer(s, footer_note or current_section["label"])
    return s


def bullets_slide(title, bullets, notes, size=19):
    s = content_slide(title)
    add_bullets(s, bullets, Inches(0.75), Inches(1.55), Inches(11.8), Inches(5.2), size=size)
    set_notes(s, notes)
    return s


def image_slide(title, image_path, notes, caption=None, max_h=Inches(5.4)):
    s = content_slide(title)
    top = Inches(1.5)
    add_picture_contain(s, image_path, Inches(0.5), top, Inches(12.3), max_h)
    if caption:
        add_text(s, caption, Inches(0.75), Inches(1.5) + max_h + Inches(0.05), Inches(11.8), Inches(0.4),
                  size=13, color=TEXT_SECONDARY, italic=True, align=PP_ALIGN.CENTER)
    set_notes(s, notes)
    return s


def image_with_bullets_slide(title, image_path, bullets, notes, img_w=Inches(6.3)):
    s = content_slide(title)
    add_picture_contain(s, image_path, Inches(0.5), Inches(1.55), img_w, Inches(5.3), align="left")
    add_bullets(s, bullets, Inches(7.1), Inches(1.7), Inches(5.6), Inches(5.0), size=16.5)
    set_notes(s, notes)
    return s


def big_stat_slide(title, stats, notes):
    """stats: list of (number, label, color)"""
    s = content_slide(title)
    n = len(stats)
    col_w = Inches(11.8) / n
    for i, (number, label, color) in enumerate(stats):
        left = Inches(0.75) + col_w * i
        add_text(s, number, left, Inches(2.6), col_w, Inches(1.3), size=54, bold=True,
                  color=color, align=PP_ALIGN.CENTER)
        add_text(s, label, left, Inches(3.9), col_w, Inches(1.0), size=15,
                  color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
    set_notes(s, notes)
    return s


def table_slide(title, headers, rows, notes, col_widths=None, font_size=13, footer_note=None):
    s = content_slide(title, footer_note=footer_note)
    n_rows, n_cols = len(rows) + 1, len(headers)
    left, top, width, height = Inches(0.6), Inches(1.55), Inches(12.1), Inches(5.3)
    table_shape = s.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(font_size + 1)
        p.font.color.rgb = WHITE
        p.font.name = FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(font_size)
            p.font.color.rgb = TEXT_PRIMARY
            p.font.name = FONT
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE if i % 2 == 0 else C("#f2f1ee")
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
    set_notes(s, notes)
    return s


def box_flow_slide(title, boxes, notes, arrow=True):
    """boxes: list of (label, sublabel, color)."""
    s = content_slide(title)
    n = len(boxes)
    gap = Inches(0.35)
    box_w = (Inches(12.2) - gap * (n - 1)) / n
    box_h = Inches(2.2)
    top = Inches(2.6)
    left = Inches(0.6)
    for i, (label, sublabel, color) in enumerate(boxes):
        x = left + i * (box_w + gap)
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        shp.shadow.inherit = False
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(17)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT
        if sublabel:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = sublabel
            run2.font.size = Pt(11.5)
            run2.font.color.rgb = C("#f0efec")
            run2.font.name = FONT
        if arrow and i < n - 1:
            ax = x + box_w + Inches(0.03)
            arrow_shp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, top + box_h / 2 - Inches(0.13),
                                            gap - Inches(0.06), Inches(0.26))
            arrow_shp.fill.solid()
            arrow_shp.fill.fore_color.rgb = MUTED
            arrow_shp.line.fill.background()
            arrow_shp.shadow.inherit = False
    set_notes(s, notes)
    return s


# ===========================================================================
# BUILD THE DECK
# ===========================================================================

# --- Title ---
title_slide(
    "Computer Vision: A Crash Course",
    "From pixels to predictions — building a real fabric defect detector",
    "Postgraduate seminar · No prior computer vision background assumed",
)
set_notes(prs.slides[-1],
    "Welcome. Today is a crash course in computer vision — the field of teaching "
    "computers to make sense of images. We will not just cover theory: everything "
    "we discuss will be demonstrated using a real project, built end to end, that "
    "classifies photos of fabric into 9 categories — 8 kinds of manufacturing "
    "defects plus 'defect-free'. By the end you should understand not just what "
    "the words 'CNN', 'transfer learning', and 'confusion matrix' mean, but how "
    "they were actually used to get a working system.")

# --- Agenda ---
s = content_slide("What We'll Cover", footer_note="Roadmap")
add_bullets(s, [
    "Part 1 — What is computer vision? How does a photo become numbers a computer can use?",
    "Part 2 — How machines learn to see: neural networks and convolutional neural networks (CNNs)",
    "Part 3 — The practical trick that makes this possible with limited data: transfer learning",
    "Part 4 — Training a model, step by step: data, loss, optimizer, epochs",
    "Part 5 — Evaluating a model properly — and reading a confusion matrix in depth",
    "Part 6 — Case study: our Fabric Defect Detector, from raw dataset to a working app",
], Inches(0.75), Inches(1.6), Inches(11.8), Inches(5.2), size=19, space_after=16)
set_notes(prs.slides[-1],
    "This is a seminar-style crash course, meaning we move quickly across the full "
    "pipeline rather than spending weeks on any one topic. The structure mirrors "
    "how a real computer vision project actually gets built: first understand the "
    "data and the problem, then understand the model architecture, then understand "
    "how training works mechanically, then understand how to judge whether it "
    "worked, and finally we walk through our own project end to end where every "
    "single concept gets grounded in a concrete, real result.")

# ===================== PART 1 =====================
section_slide(1, "What Is Computer Vision?",
              "Teaching a computer to extract meaning from images and video")

s = content_slide("What Is Computer Vision?")
add_bullets(s, [
    "A field of AI focused on getting computers to understand visual information — photos, video, camera feeds",
    "“Understand” can mean many things: what object is this? where is it? what's its boundary? is something wrong with it?",
    "Powers things you already use: face unlock on your phone, self-driving car lane detection, medical scan analysis, quality inspection on factory lines",
    "Our case study today is that last one: automatically spotting defects in fabric — a real textile manufacturing quality-control problem",
], Inches(0.75), Inches(1.55), Inches(11.8), Inches(5.2), size=19, space_after=14)
set_notes(prs.slides[-1],
    "Computer vision is the umbrella field for any AI system that takes pixels as "
    "input and produces some kind of understanding as output. It is a subfield of "
    "machine learning, specialized around visual data. The reason it gets its own "
    "field, rather than just being 'AI applied to images', is that images have a "
    "particular structure — spatial structure, where nearby pixels are related — "
    "and that structure demands specialized techniques, which is exactly what "
    "we'll build up to with convolutional neural networks. Our running example: "
    "textile factories currently rely on human inspectors staring at fabric rolls "
    "looking for holes, stains, broken stitches. That's slow and inconsistent "
    "between inspectors. The question this project asks is whether a model can "
    "do this instead, from a photograph.")

image_slide(
    "How Images Become Numbers",
    os.path.join(ASSETS, "pixel_grid.png"),
    "A computer has no concept of 'green fabric' or 'a hole' — it only ever sees "
    "numbers. Every image is stored as a grid of pixels (short for 'picture "
    "elements'), and every pixel is just a brightness number, typically 0 to 255, "
    "where 0 is black and 255 is white. Zoom into any photo far enough and this "
    "is all there is — no shapes, no objects, just a big grid of numbers. This "
    "matters because it reframes the entire problem: computer vision is really "
    "the problem of finding mathematical operations on grids of numbers that "
    "happen to correlate with the visual concepts we care about, like 'this "
    "looks like a hole'.",
    caption="A 224×224 pixel image is a grid of 50,176 brightness numbers — shown zoomed in as an 8×8 patch",
)

image_slide(
    "Color Images = 3 Stacked Grids",
    os.path.join(ASSETS, "rgb_channels.png"),
    "A grayscale image is one grid of numbers. A color image is three grids "
    "stacked together — one for red intensity, one for green, one for blue — "
    "called channels. Mix the three intensities per pixel and you get any color. "
    "So a standard color photo resized to 224×224 pixels is really a block of "
    "224 × 224 × 3 = about 150,000 numbers. That block of numbers, called a "
    "tensor, is the actual input our model receives — there is no image in the "
    "traditional sense inside the computer, only this tensor. Every operation "
    "we'll discuss (convolutions, layers, etc.) is just math performed on this "
    "block of numbers.",
    caption="Notice the faint vertical mark (a real hole defect) is visible, though faintly, in all three channels",
)

s = box_flow_slide(
    "The Three Core Computer Vision Tasks",
    [
        ("Classification", "“What is this, overall?”\nOne label for the whole image\n→ our project", BLUE),
        ("Detection", "“What's here, and where?”\nBoxes around each object", VIOLET),
        ("Segmentation", "“Exactly which pixels?”\nA precise outline, pixel by pixel", MAGENTA),
    ],
    "These three tasks are progressively more detailed answers to the same "
    "underlying question. Classification just assigns one label to the whole "
    "photo — that's exactly what our fabric defect detector does: given a whole "
    "image, say 'hole' or 'stain' or 'defect-free'. Detection goes further and "
    "draws a bounding box around each object of interest — useful if there could "
    "be multiple defects in one photo and you need to know where each one is. "
    "Segmentation goes furthest and traces the exact pixel-level outline of the "
    "defect, not just a box. We chose classification for this project because "
    "the dataset only had one label per image, and it's the simplest starting "
    "point; the README explicitly lists detection (with a technique called "
    "YOLO) and segmentation as future improvements once more granular data "
    "is available.",
)

# ===================== PART 2 =====================
section_slide(2, "How Machines Learn to See",
              "Neural networks, convolutions, and the CNN feature hierarchy")

s = content_slide("Neural Networks in 60 Seconds")
add_bullets(s, [
    "A neural network is a stack of layers, each made of many simple units (“neurons”)",
    "Each neuron takes numbers in, multiplies them by learned weights, adds them up, and passes the result through a small nonlinear function",
    "Stack enough of these layers and the network can approximate very complex functions — including “is this a picture of a hole?”",
    "“Learning” = adjusting all those weights, gradually, based on many labeled examples — nobody hand-codes the rules",
], Inches(0.75), Inches(1.55), Inches(11.8), Inches(4.6), size=19, space_after=14)
set_notes(prs.slides[-1],
    "You don't need the full mathematics to follow the rest of this course, just "
    "the shape of the idea. A neuron computes a weighted sum of its inputs, plus "
    "a small twist called an activation function that lets the network represent "
    "non-linear relationships — without it, stacking layers would collapse into "
    "one big linear equation, which can't represent something as complex as "
    "'contains a hole'. What makes neural networks powerful is that we don't "
    "design the weights by hand; we initialize them randomly and then let an "
    "algorithm called backpropagation adjust them automatically based on "
    "thousands of labeled examples, which we'll cover in Part 4.")

s = content_slide("Why Not Just Feed Pixels to a Plain Neural Network?")
add_bullets(s, [
    "A 224×224×3 image is ~150,000 numbers — a plain (“fully connected”) network would need a separate weight for every input-to-neuron connection",
    "That's millions of weights before the network has learned anything useful — slow, prone to overfitting, and it ignores spatial structure entirely",
    "Spatial structure matters: a hole looks like a hole regardless of where in the photo it appears — a plain network has to re-learn that separately for every position",
    "Convolutional Neural Networks (CNNs) solve both problems at once, using a small building block called convolution",
], Inches(0.75), Inches(1.55), Inches(11.8), Inches(5.0), size=18.5, space_after=13)
set_notes(prs.slides[-1],
    "This slide exists to motivate why CNNs were invented rather than just using "
    "generic neural networks. Two separate problems compound: raw scale (too "
    "many weights to learn from limited data) and a missing assumption (images "
    "have local structure — nearby pixels are correlated, and patterns are "
    "position-independent, meaning a hole in the top-left of a photo looks the "
    "same as a hole in the bottom-right). Convolution bakes both efficiencies in "
    "by design, which is the single biggest reason CNNs became the default tool "
    "for vision starting in the early 2010s.")

image_slide(
    "Convolution: The Core Building Block of a CNN",
    os.path.join(ASSETS, "convolution.png"),
    "A convolution filter (or kernel) is a small grid of numbers, say 3×3, that "
    "represents a pattern to look for — an edge, a corner, a texture. The filter "
    "slides across the whole image, and at every position it multiplies its "
    "numbers against the pixels underneath and sums the result into one output "
    "number. High output = strong match at that location. Do this everywhere and "
    "you get a new grid called a feature map, showing where that pattern occurs "
    "across the whole image. Critically, the same filter (the same small set of "
    "numbers) is reused at every position — this is called weight sharing, and "
    "it's what keeps the number of learnable parameters small regardless of "
    "image size, while naturally handling 'a hole looks like a hole anywhere in "
    "the photo'.",
)

image_slide(
    "Stacking Convolutions: A Hierarchy of Features",
    os.path.join(ASSETS, "feature_hierarchy.png"),
    "A CNN doesn't use one filter — it stacks many layers of many filters. This "
    "is a real example: we fed one of our actual fabric photos through our "
    "trained network and extracted what an early layer 'sees' versus a late "
    "layer. Early layers detect simple, generic patterns — edges, corners, "
    "fabric weave texture — the same kind of features that are useful for any "
    "image, not just fabric. Later layers combine those simple detections into "
    "increasingly abstract, complex patterns — shapes specific to the objects in "
    "the training data. This progression, simple to complex, is why CNNs are "
    "described as learning a feature hierarchy automatically, rather than "
    "requiring a human to hand-design 'what a hole looks like'.",
    caption="Real activations extracted from our trained ResNet18, layer1 (early) vs layer4 (late), on an actual hole-defect photo",
)

s = content_slide("Popular CNN Architectures (the ones we used)")
add_bullets(s, [
    ("ResNet18 — the backbone we trained with", 0),
    ("Introduced “residual connections”: shortcuts that let signal skip layers, making very deep networks trainable without degrading", 1),
    ("“18” = 18 layers deep; small and fast enough to run on a CPU, which mattered for this project", 1),
    ("MobileNetV2 — the lightweight alternative also supported in our code", 0),
    ("Designed for phones/edge devices — fewer parameters, faster inference, small accuracy trade-off", 1),
    ("Both were pretrained on ImageNet before we ever touched them — which is the subject of Part 3", 0),
], Inches(0.75), Inches(1.55), Inches(11.8), Inches(5.2), size=18, space_after=12)
set_notes(prs.slides[-1],
    "You don't need to memorize architecture internals for this course, but it's "
    "worth knowing these names because they recur constantly in real projects. "
    "ResNet (Residual Network) solved a specific historical problem: before it, "
    "researchers found that simply stacking more convolutional layers made "
    "networks perform worse, not better, past a certain depth, because gradients "
    "shrink as they propagate backward through many layers ('vanishing "
    "gradients'). Residual connections — shortcuts that let a layer's input skip "
    "ahead and be added to its output — fixed this and enabled genuinely deep "
    "networks. We picked ResNet18, the smallest ResNet variant, deliberately: it "
    "is accurate enough for this problem while still being feasible to run "
    "without a GPU. MobileNetV2 is included in our code as a lighter-weight "
    "alternative, more relevant if this model needed to run on a factory-floor "
    "device rather than a server.")

# ===================== PART 3 =====================
section_slide(3, "The Practical Trick: Transfer Learning",
              "How to get a strong model with a dataset of only ~3,000 images")

s = content_slide("The Data Problem")
add_bullets(s, [
    "Training a CNN from scratch (“random weights”) typically needs hundreds of thousands to millions of labeled images to work well",
    "Our fabric dataset has 3,067 images across 9 classes — nowhere near enough to train a CNN from zero",
    "This is the normal situation for real-world projects: labeled data is expensive and slow to collect, especially for a niche domain like industrial fabric defects",
    "The solution the field converged on: don't start from zero — start from a model that already understands images in general",
], Inches(0.75), Inches(1.55), Inches(11.8), Inches(5.0), size=18.5, space_after=13)
set_notes(prs.slides[-1],
    "This is the single most important practical lesson of this course, because "
    "it's the difference between a project being feasible or not. Most real "
    "organizations do not have millions of labeled images for their specific "
    "problem — they have a few hundred to a few thousand, exactly like us. If "
    "the field required training every CNN from scratch, computer vision would "
    "be inaccessible outside a handful of large tech companies. Transfer "
    "learning is the technique that changed that.")

s = content_slide("Enter ImageNet & Pretrained Models")
add_bullets(s, [
    "ImageNet: a public dataset of ~1.4 million photos across 1,000 general categories — dogs, cars, furniture, fruit, and so on",
    "Researchers train large CNNs (like ResNet18) on ImageNet once, then release the trained weights publicly — a “pretrained model”",
    "Those weights already encode a lot of general visual knowledge: how to detect edges, textures, shapes, object parts",
    "That knowledge turns out to transfer surprisingly well to completely different domains — including fabric photographs, which look nothing like ImageNet's photos",
], Inches(0.75), Inches(1.55), Inches(11.8), Inches(5.0), size=18.5, space_after=13)
set_notes(prs.slides[-1],
    "It's worth pausing on why this transfer works at all, since ImageNet has no "
    "fabric-defect photos in it. The answer is that the earliest layers of a CNN "
    "learn extremely generic visual primitives — edge orientation, color "
    "gradients, simple textures — and those primitives are useful for "
    "essentially any visual task, not just recognizing dogs. It's only the "
    "later, more abstract layers that become specific to the original training "
    "categories. That's exactly the structure we saw two slides ago in the "
    "feature hierarchy image.")

s = new_slide()
fill_bg(s, SURFACE)
add_title_bar(s, "Transfer Learning: Reuse the Backbone, Train a New Head")
add_footer(s, current_section["label"])
box_top = Inches(2.0)
add_rect(s, Inches(0.8), box_top, Inches(5.6), Inches(2.6), BLUE)
add_text(s, "BACKBONE", Inches(0.8), box_top + Inches(0.25), Inches(5.6), Inches(0.5),
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Pretrained on ImageNet\nAlready knows: edges, textures,\nshapes, object parts",
         Inches(0.8), box_top + Inches(0.85), Inches(5.6), Inches(1.5),
         size=15, color=C("#cde2fb"), align=PP_ALIGN.CENTER, line_spacing=1.3)
arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.55), box_top + Inches(1.0), Inches(0.7), Inches(0.6))
arrow.fill.solid(); arrow.fill.fore_color.rgb = MUTED; arrow.line.fill.background(); arrow.shadow.inherit = False
add_rect(s, Inches(7.4), box_top, Inches(5.1), Inches(2.6), ORANGE)
add_text(s, "HEAD (new)", Inches(7.4), box_top + Inches(0.25), Inches(5.1), Inches(0.5),
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Randomly initialized\nTrained from scratch on OUR\n3,067 fabric images → 9 classes",
         Inches(7.4), box_top + Inches(0.85), Inches(5.1), Inches(1.5),
         size=15, color=C("#fde3d3"), align=PP_ALIGN.CENTER, line_spacing=1.3)
add_text(s, "This is transfer learning: keep the general-purpose visual knowledge, only teach the model our specific 9 categories.",
         Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.6), size=16, italic=True, color=TEXT_SECONDARY,
         align=PP_ALIGN.CENTER)
set_notes(s,
    "This diagram is the mental model to hold onto for the rest of the course. "
    "Every CNN we use has two conceptual parts: a backbone — all the "
    "convolutional layers that extract visual features — and a head — the final "
    "layer(s) that map those features to actual class predictions. In transfer "
    "learning, we take a backbone that arrives already trained (on ImageNet) and "
    "we throw away the original head, because it was built for ImageNet's 1,000 "
    "categories, not our 9. We attach a brand-new head, randomly initialized, "
    "sized for our 9 fabric classes, and that's the only part that starts "
    "learning from zero.")

s = content_slide("Frozen vs. Fine-Tuned: The Trade-off We Actually Made")
add_bullets(s, [
    ("Option A — Freeze the backbone (what we ran locally, on CPU)", 0),
    ("Backbone weights are locked; only the new head trains", 1),
    ("Much faster, works fine with limited data and no GPU — but the visual features stay generic, never adapted to fabric specifically", 1),
    ("Option B — Fine-tune the whole network (what our Colab notebook does, on GPU)", 0),
    ("Backbone weights are allowed to adjust too, usually with a small learning rate", 1),
    ("Slower and needs more compute, but lets early layers specialize — e.g. learning to tell fabric weave directions apart", 1),
], Inches(0.75), Inches(1.55), Inches(11.8), Inches(5.2), size=18, space_after=11)
set_notes(prs.slides[-1],
    "This is a direct, real engineering decision we made in this project, not a "
    "hypothetical. We had no GPU on the local machine, so we froze the backbone "
    "entirely — meaning during training, gradients were only computed for the "
    "final linear layer, which made each epoch dramatically faster since no "
    "backpropagation had to run through the convolutional layers at all. This "
    "got us a working model, 90% validation accuracy, in under 40 minutes on a "
    "CPU. The trade-off, which we'll see concretely in the results later, is "
    "that some fabric-specific visual distinctions — particularly telling apart "
    "different orientations of line-shaped defects — were never something the "
    "ImageNet backbone needed to learn, and freezing it means it never got the "
    "chance to adapt. The project's Colab notebook is deliberately set up to "
    "fine-tune the whole network on a free GPU, precisely to close that gap.")

# ===================== PART 4 =====================
section_slide(4, "Training a Model, Step by Step",
              "Preprocessing, splitting, loss, optimization, and the training loop")

s = content_slide("Step 1: Preprocessing & Data Augmentation")
add_bullets(s, [
    "Preprocessing: every image is resized to a fixed size (224×224) and its pixel values are rescaled to match what the pretrained backbone expects",
    "Data augmentation: during training only, images are randomly flipped, rotated, and color-jittered",
    "Why augment? With only ~2,000 training images, a model can memorize specifics instead of learning general patterns (overfitting) — augmentation manufactures variety for free",
    "A rotated, slightly recolored photo of a hole is still a hole — augmentation teaches the model that invariance directly",
], Inches(0.75), Inches(1.55), Inches(11.8), Inches(5.0), size=18.5, space_after=13)
set_notes(prs.slides[-1],
    "Preprocessing and augmentation are easy to underestimate but matter enormously "
    "with small datasets like ours. Resizing to a fixed resolution is mechanical "
    "— the network architecture expects a fixed-size input tensor. Normalization "
    "(rescaling pixel values) matters because the ImageNet-pretrained backbone was "
    "trained on images normalized a particular way, and if we feed it images on "
    "a different numeric scale, its pretrained weights won't behave as expected. "
    "Augmentation is the more interesting idea: since we can't collect more real "
    "photos, we synthesize plausible variations of the ones we have, which acts "
    "as a regularizer — it makes it harder for the model to just memorize the "
    "training set and forces it to learn features that generalize.")

image_slide(
    "Step 2: Splitting the Data — Train / Validation / Test",
    os.path.join(ASSETS, "split_chart.png"),
    "We never train and evaluate on the same images — that would let the model "
    "'cheat' by memorizing answers instead of learning generalizable patterns. So "
    "the dataset gets split three ways: the training set is what the model "
    "actually learns from; the validation set is checked after every epoch to "
    "monitor progress and decide when to stop, without ever being trained on; "
    "the test set is touched exactly once, at the very end, to report an honest, "
    "final accuracy number. Our split follows a standard 70/15/15 ratio.",
)

s = content_slide("Step 3: The Loss Function")
add_bullets(s, [
    "After each prediction, we need a single number that says “how wrong was that?” — that's the loss",
    "We use cross-entropy loss, standard for classification: the model outputs a probability for each of the 9 classes",
    "Cross-entropy penalizes confident wrong answers much more harshly than unsure wrong answers",
    "Example: predicting 99% “stain” when the truth is “hole” gets a much bigger penalty than predicting 40% “stain” for the same wrong guess",
    "Training = repeatedly trying to make this one number smaller, across the whole training set",
], Inches(0.75), Inches(1.55), Inches(11.8), Inches(5.2), size=18, space_after=12)
set_notes(prs.slides[-1],
    "The loss function is the single quantity that training is actually trying "
    "to minimize — everything else in the pipeline exists to compute it and "
    "reduce it. Cross-entropy loss comes from information theory and has a "
    "specific mathematical property that makes it ideal for classification: its "
    "penalty grows sharply as a confident prediction moves further from the "
    "truth. This shapes the model's behavior during training — it's actively "
    "discouraged from being both wrong and overconfident, which is exactly the "
    "behavior we want from a system whose confidence scores humans might trust.")

s = content_slide("Step 4: The Optimizer")
add_bullets(s, [
    "Knowing “how wrong” isn't enough — we need to know which direction to nudge every single weight to make it less wrong",
    "Backpropagation computes that direction for every weight in the network (the “gradient”), using calculus applied automatically",
    "Gradient descent: take a small step in the direction that reduces the loss — repeat, over and over",
    "We used Adam, a popular refinement of gradient descent that adapts the step size per weight automatically, converging faster and more reliably than plain gradient descent",
], Inches(0.75), Inches(1.55), Inches(11.8), Inches(5.2), size=18, space_after=12)
set_notes(prs.slides[-1],
    "Think of the loss as a landscape with hills and valleys, where the height at "
    "any point represents how wrong the model currently is for a particular "
    "setting of all its weights. Gradient descent is the strategy of always "
    "stepping downhill — backpropagation is the mechanism that tells you which "
    "direction is downhill from wherever you currently stand, for every one of "
    "potentially millions of weights simultaneously. Adam improves on the basic "
    "version by keeping a running memory of recent gradients per weight, which "
    "lets it take bigger, more confident steps in consistent directions and "
    "smaller steps where the landscape is noisy — in practice this converges "
    "faster and is far more forgiving of the learning-rate choice, which is why "
    "it's the default optimizer in most modern deep learning code, including "
    "ours.")

s = content_slide("Step 5: Epochs & the Training Loop")
add_bullets(s, [
    "One epoch = one full pass through every training image, once",
    ("For every batch of images in an epoch:", 0),
    ("Forward pass — run the images through the network to get predictions", 1),
    ("Compute the loss — how wrong were those predictions?", 1),
    ("Backward pass — backpropagation computes the gradient for every trainable weight", 1),
    ("Optimizer step — Adam nudges every weight slightly, using that gradient", 1),
    ("Repeat for many epochs, checking validation accuracy after each one, until it stops improving (“early stopping”)", 0),
], Inches(0.75), Inches(1.5), Inches(11.8), Inches(5.3), size=17, space_after=9)
set_notes(prs.slides[-1],
    "This is the complete training loop, and every deep learning training run in "
    "existence is some variation of this exact sequence, repeated thousands of "
    "times. Batching — processing a handful of images at a time (we used 32) "
    "rather than one at a time or the whole dataset at once — exists purely as a "
    "practical compromise between memory usage and how noisy versus stable each "
    "weight update is. Early stopping is a safeguard against overfitting: we "
    "watch validation accuracy, and if it stops improving for a set number of "
    "epochs (we used a patience of 4), we stop, on the logic that further "
    "training would just be memorizing the training set rather than genuinely "
    "improving.")

# ===================== PART 5 =====================
section_slide(5, "Evaluating a Model",
              "Why accuracy alone can mislead you, and how to read a confusion matrix")

s = content_slide("Why “Accuracy” Alone Can Lie to You")
add_bullets(s, [
    "Accuracy = fraction of predictions that were correct — simple, but dangerous with imbalanced data",
    "Our dataset is imbalanced: “defect free” is 260 of 460 test images (~57%), while “Vertical” defects are only 14",
    "A lazy model that always predicts “defect free”, regardless of the image, would score ~57% accuracy — while being completely useless",
    "We need metrics that account for performance on each individual class, not just the overall average",
], Inches(0.75), Inches(1.55), Inches(11.8), Inches(5.0), size=18.5, space_after=13)
set_notes(prs.slides[-1],
    "This is a crucial, commonly-missed lesson for anyone starting out in applied "
    "machine learning. Accuracy is intuitive and easy to report, which is exactly "
    "why it's so often misused. Whenever a dataset is imbalanced — and most "
    "real-world datasets are, ours very much included — a high accuracy number "
    "can be hiding a model that has essentially learned to always guess the "
    "majority class. This is precisely why our evaluation script reports a full "
    "classification report and a confusion matrix, not just one accuracy "
    "percentage.")

s = new_slide()
fill_bg(s, SURFACE)
add_title_bar(s, "Precision, Recall, and F1 — Defined Simply")
add_footer(s, current_section["label"])
# 2x2 quadrant
qx, qy, qw, qh = Inches(1.0), Inches(1.7), Inches(5.6), Inches(3.6)
labels = [
    ("Predicted: hole\nActually: hole", "True Positive", GREEN),
    ("Predicted: hole\nActually: not hole", "False Positive", RED),
    ("Predicted: not hole\nActually: hole", "False Negative", ORANGE),
    ("Predicted: not hole\nActually: not hole", "True Negative", MUTED),
]
positions = [(qx, qy), (qx + qw / 2, qy), (qx, qy + qh / 2), (qx + qw / 2, qy + qh / 2)]
for (sub, tag, color), (x, y) in zip(labels, positions):
    add_rect(s, x, y, qw / 2 - Inches(0.05), qh / 2 - Inches(0.05), color)
    add_text(s, tag, x + Inches(0.15), y + Inches(0.12), qw / 2 - Inches(0.3), Inches(0.4),
              size=13, bold=True, color=WHITE)
    add_text(s, sub, x + Inches(0.15), y + Inches(0.55), qw / 2 - Inches(0.3), Inches(1.4),
              size=11.5, color=WHITE, line_spacing=1.2)
add_bullets(s, [
    ("Precision = of everything predicted “hole”, what fraction really was?", 0),
    ("TP / (TP + FP) — asks: can I trust a positive prediction?", 1),
    ("Recall = of everything that really was “hole”, what fraction did we catch?", 0),
    ("TP / (TP + FN) — asks: do I miss real cases?", 1),
    ("F1 = a single number balancing precision and recall (their harmonic mean)", 0),
    ("Useful when you need one number but don't want to ignore either type of mistake", 1),
], Inches(7.0), Inches(1.75), Inches(5.6), Inches(4.8), size=15.5, space_after=12)
set_notes(s,
    "Precision and recall are the two lenses that a single accuracy number "
    "collapses into one, and they answer different practical questions. "
    "Precision matters when false positives are costly — e.g. if flagging a "
    "perfectly good piece of fabric as defective causes it to be needlessly "
    "discarded, you care about precision. Recall matters when false negatives "
    "are costly — e.g. if a real hole slipping through to customers is worse "
    "than an occasional false alarm, you care about recall. F1 is a compromise "
    "metric when you need to rank models with one number but don't want to "
    "ignore either failure mode. In our results, we'll see these numbers differ "
    "quite a bit per class, which is exactly the kind of insight a single "
    "accuracy percentage would have hidden.")

image_slide(
    "The Confusion Matrix, Fully Explained",
    os.path.join(ASSETS, "confusion_annotated.png"),
    "A confusion matrix is a grid: rows are the true label, columns are the "
    "model's predicted label, and each cell counts how many test images fall "
    "into that (truth, prediction) pair. A perfect model would have every image "
    "on the diagonal (true label = predicted label). Reading across a row (like "
    "the red-outlined ‘Vertical’ row here) tells you recall — of all the real "
    "Vertical images, where did they actually get predicted? Reading down a "
    "column (like the violet-outlined ‘hole’ column) tells you precision — of "
    "everything the model called ‘hole’, how much of it came from other true "
    "classes bleeding in? This single picture contains far more diagnostic "
    "information than any single accuracy number — it tells you exactly which "
    "pairs of classes the model confuses, not just how often it's wrong.",
    max_h=Inches(5.55),
)

# ===================== PART 6 =====================
section_slide(6, "Case Study: Fabric Defect Detector",
              "Every concept from Parts 1–5, applied end to end, with real results")

s = content_slide("The Problem We Set Out to Solve")
add_bullets(s, [
    "Manual fabric inspection on a production line is slow and inconsistent between human inspectors",
    "Goal: can a lightweight, pretrained CV model classify common fabric defects from a photograph?",
    "Relevant to real textile QC and export-compliance workflows, where defective fabric must be caught before shipment",
    "Deliberately scoped as a learning project + portfolio piece — realistic constraints: small dataset, no dedicated GPU",
], Inches(0.75), Inches(1.55), Inches(11.8), Inches(5.0), size=19, space_after=14)
set_notes(prs.slides[-1],
    "Every project should start with a clear problem statement, and it's worth "
    "restating ours plainly before diving into implementation details: textile "
    "manufacturers need to catch defective fabric before it ships, and doing "
    "that by eye, continuously, is expensive and error-prone. This project asks "
    "whether an image classifier, built with realistic constraints — a modest "
    "dataset, no dedicated GPU on the development machine — can meaningfully "
    "help. That constraint framing matters: it's what justified transfer "
    "learning over training from scratch, and a frozen backbone over full "
    "fine-tuning, for the local development loop.")

s = content_slide("The Dataset")
add_bullets(s, [
    "Multi-Class Fabric Defect Detection Dataset (Kaggle) — real production-line photos from Basler industrial cameras",
    "3,067 images across 9 classes: defect free, hole, horizontal, Vertical, lines, Pinched fabric, Needle mark, Broken stitch, stain",
    "Meaningfully imbalanced: “defect free” alone is ~54% of all images; “Vertical” is only ~3%",
    "Flat folder layout (one folder per class, no pre-made train/val/test split) — our data loader auto-detects this and splits 70/15/15 itself",
], Inches(0.75), Inches(1.55), Inches(11.8), Inches(5.0), size=18.5, space_after=13)
set_notes(prs.slides[-1],
    "Understanding your dataset before writing any modeling code is a "
    "non-negotiable step in any real project, and it directly predicts the "
    "failures you'll see later. The class imbalance here — defect-free "
    "dominating, several defect types having barely 100 examples — isn't a flaw "
    "in the dataset, it's realistic: in a well-run factory, defects should be "
    "rare. But it does mean the model gets far more practice recognizing "
    "'normal' fabric than recognizing the rarest defect types, which we will see "
    "reflected directly in the confusion matrix later.")

s = content_slide("Our Architecture & Training Choices")
add_bullets(s, [
    ("Backbone: ResNet18, pretrained on ImageNet (see Part 3)", 0),
    ("Head: replaced with a new linear layer outputting 9 fabric classes", 0),
    ("Local run (this demo): backbone frozen — only the head trains, on CPU, ~200s/epoch", 0),
    ("Recommended run: notebooks/train_model.ipynb on Google Colab's free GPU — fine-tunes the whole backbone", 0),
    ("Preprocessing: resize to 224×224, ImageNet normalization, random flip/rotate/color-jitter augmentation on training data only", 0),
    ("Optimizer: Adam, learning rate 1e-3, cross-entropy loss, up to 12 epochs with early stopping", 0),
], Inches(0.75), Inches(1.5), Inches(11.8), Inches(5.3), size=17.5, space_after=12)
set_notes(prs.slides[-1],
    "This slide is a direct summary of the actual command that was run: "
    "python src/train.py with resnet18, a frozen backbone, batch size 32, up to "
    "12 epochs, patience of 4 for early stopping, and a learning rate of 1e-3 "
    "chosen higher than typical full-fine-tuning rates precisely because only "
    "the small new head layer needed to learn, which tolerates a larger step "
    "size than adjusting a whole pretrained network would.")

image_slide(
    "Training in Action — Our Actual Curves",
    os.path.join(ASSETS, "training_curves.png"),
    "This is not a textbook illustration — these are the real loss and accuracy "
    "curves from our training run, read directly from the logged history. Loss "
    "drops sharply in the first couple of epochs then levels off; accuracy rises "
    "the same way, reaching 90% on validation by epoch 12. Notice train and "
    "validation stay close together throughout, rather than train shooting far "
    "ahead of validation — that closeness is the visual signature of a model "
    "that is not overfitting. Training ran the full 12-epoch budget without "
    "triggering early stopping, meaning validation accuracy was still "
    "(slowly) improving when we stopped.",
)

big_stat_slide(
    "Headline Results",
    [
        ("90.0%", "Validation accuracy", BLUE),
        ("88.5%", "Test accuracy\n(460 held-out images)", AQUA),
        ("8", "Classes with F1 > 0.75\nout of 9 total", VIOLET),
    ],
    "These are the three numbers to remember. Validation accuracy (90%) is what "
    "we monitored during training to pick the best checkpoint. Test accuracy "
    "(88.5%) is the honest, final number — measured once, on the 460 test "
    "images the model never saw in any form during training or checkpoint "
    "selection. The small gap between the two (90% vs 88.5%) is expected and "
    "healthy — a large gap would suggest we'd tuned decisions based on the "
    "validation set until we were effectively overfitting to it. Per-class F1 "
    "scores ranged from 0.48 (Vertical, our weakest class) to 0.97 (defect "
    "free, our strongest) — which is exactly what the confusion matrix on the "
    "next slide explains.",
)

s = image_slide(
    "Reading Our Own Confusion Matrix",
    os.path.join(MODELS, "confusion_matrix.png"),
    "Two very different stories live in this matrix. Story one: defect free "
    "(247/260 correct) and stain (52/56) are excellent — both have a distinctive "
    "visual signature (a blotchy texture, or simply the absence of any mark) "
    "that the frozen, general-purpose ImageNet features already captured well. "
    "Story two: hole, horizontal, Vertical, and lines form a confused cluster — "
    "Vertical is right only 5 of 14 times, frequently mistaken for hole or "
    "horizontal. Why this specific cluster and not others? These four classes "
    "are all fundamentally 'a linear or gap-shaped irregularity running in some "
    "direction across the weave' — they differ mainly by orientation and "
    "thickness, a distinction ImageNet's original 1,000 categories (dogs, cars, "
    "furniture) never needed to make, and one our frozen backbone was therefore "
    "never forced to learn. This is the single clearest illustration in the "
    "whole project of the frozen-backbone trade-off discussed in Part 3 — and "
    "precisely the gap that fine-tuning the full network (our Colab notebook) "
    "is positioned to close.",
    max_h=Inches(5.55),
)

image_slide(
    "From Model to Product: The Streamlit Demo",
    os.path.join(ASSETS, "app_mockup.png"),
    "A trained model sitting in a file is not useful to anyone by itself — "
    "deployment is the step that turns it into a product. We wrapped the model "
    "in a small Streamlit web app: a user uploads a photo, the app runs the "
    "exact same preprocessing used during training, passes it through the "
    "model, and displays the predicted class with a confidence score for every "
    "class. This is a real result from this app, on an actual dataset image: "
    "the model's top guess was 'hole' at 60% confidence, with 'horizontal' and "
    "'Vertical' as its next guesses — exactly the confusion cluster we just "
    "identified in the confusion matrix, now visible live in the product "
    "itself, not just in an offline evaluation report.",
    max_h=Inches(5.1),
)

s = content_slide("What We'd Do Next")
add_bullets(s, [
    ("Fine-tune the full backbone on a GPU (the Colab notebook is already built for this)", 0),
    ("Should directly target the hole / horizontal / Vertical / lines confusion cluster", 1),
    ("Object detection (YOLO) — localize where the defect is, not just classify the whole image", 0),
    ("Segmentation — trace the exact pixel boundary of a defect for more precise QC decisions", 0),
    ("More data for the rarest classes (Vertical, Needle mark, Pinched fabric each had only ~13–14 test images)", 0),
    ("Deploy as a hosted web app rather than a locally-run Streamlit instance", 0),
], Inches(0.75), Inches(1.55), Inches(11.8), Inches(5.2), size=18, space_after=13)
set_notes(prs.slides[-1],
    "Every real project ends with a clear-eyed list of limitations and next "
    "steps rather than a claim of being finished, and this project is no "
    "exception. The single highest-leverage next step, directly motivated by "
    "the confusion matrix we just read, is fine-tuning the full backbone — "
    "the infrastructure for this already exists in the Colab notebook. Beyond "
    "that, object detection and segmentation represent moving from Part 1's "
    "'classification' task to the more detailed 'detection' and 'segmentation' "
    "tasks — useful once the business need shifts from 'is this fabric bad?' "
    "to 'exactly where, and how large, is the defect?'.")

# ===================== WRAP-UP =====================
s = content_slide("Key Takeaways", footer_note="Wrap-up")
add_bullets(s, [
    "Images are just grids of numbers — computer vision is finding math that correlates with visual meaning",
    "CNNs use convolution (small, reusable filters) to build a hierarchy from simple edges to complex shapes",
    "Transfer learning — reusing an ImageNet-pretrained backbone — makes strong models possible with small datasets",
    "Freezing vs. fine-tuning the backbone is a real, concrete speed/accuracy trade-off, not just theory",
    "Accuracy alone can mislead on imbalanced data — precision, recall, F1, and the confusion matrix tell the fuller story",
    "A confusion matrix's rows show recall, its columns show precision, and its off-diagonal cells name exact confusions",
    "A model is only useful once deployed — wrapping it in an app is part of the job, not an afterthought",
], Inches(0.75), Inches(1.5), Inches(11.8), Inches(5.3), size=17.5, space_after=13)
set_notes(prs.slides[-1],
    "If only one slide survives from this deck, it should be this one. Every "
    "bullet here maps to a concrete moment in the case study we just walked "
    "through — encourage students to connect each one back to a specific slide "
    "or result rather than treating them as abstract statements.")

table_slide(
    "Glossary — Quick Reference",
    ["Term", "Plain-language definition"],
    [
        ["Pixel", "One brightness/color number in the grid that makes up a digital image"],
        ["Channel", "One color layer of an image (red, green, or blue); a color image = 3 stacked channels"],
        ["CNN", "Convolutional Neural Network — a network built from convolution filters, suited to images"],
        ["Convolution / filter", "A small grid of numbers that slides across an image detecting one pattern"],
        ["Feature map", "The output grid produced by one filter, showing where its pattern occurs"],
        ["Backbone", "The convolutional layers of a CNN that extract visual features"],
        ["Head", "The final layer(s) that turn extracted features into class predictions"],
        ["Transfer learning", "Reusing a model pretrained on one dataset (ImageNet) for a new task"],
        ["Fine-tuning", "Continuing to train (some or all of) a pretrained model's weights on new data"],
        ["Epoch", "One complete pass through the entire training dataset"],
        ["Loss function", "A single number measuring how wrong the model's predictions were"],
        ["Optimizer (Adam)", "The algorithm that adjusts weights to reduce the loss, step by step"],
        ["Overfitting", "Memorizing training data instead of learning patterns that generalize"],
        ["Precision / Recall", "Precision: are positive predictions trustworthy? Recall: are real cases caught?"],
        ["Confusion matrix", "A grid of true label × predicted label counts, revealing exact model mistakes"],
    ],
    "This glossary is meant as a leave-behind reference, not something to read "
    "aloud line by line — point students to it for review before an assessment.",
    col_widths=[Inches(2.8), Inches(9.3)],
    font_size=13,
    footer_note="Wrap-up",
)

s = new_slide()
fill_bg(s, TEXT_PRIMARY)
add_text(s, "Questions?", Inches(0.9), Inches(2.9), Inches(11), Inches(1.2), size=44, bold=True, color=WHITE)
add_text(s, "Code, notebook, trained model, and this deck: all in the project repository.",
         Inches(0.9), Inches(4.1), Inches(11), Inches(0.6), size=17, color=C("#c3c2b7"))
set_notes(s, "Open the floor for questions. Good prompts if the room is quiet: 'which "
             "part of the confusion matrix surprised you?', or 'where would you reach "
             "for detection or segmentation instead of classification in your own "
             "field?'.")

prs.save(OUT_PATH)
print(f"Saved {OUT_PATH} with {len(prs.slides)} slides")
